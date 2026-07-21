"""
edit_factsheet_junio.py
=======================
Actualiza Factsheet BIG Junio 2026 preservando estetica del template Mayo 2026.
Approach: replace_in_runs shape por shape (indices fijos), NUNCA borrar runs.
"""
import copy
import sys
from pathlib import Path
from pptx import Presentation

INPUT = Path("C:/Users/lmonp/Dropbox/Banca Privada (1)/AMC PAMPA CAPITAL/BIG Factsheets/2026/Pampa Capital AM - Factsheet Junio26.pptx")


def replace_run_text(shape, old_text: str, new_text: str) -> bool:
    """Reemplaza texto en el shape preservando runs y estilos.
    Busca el old_text en cualquier run individual y lo reemplaza.
    Devuelve True si hizo un reemplazo."""
    if not shape.has_text_frame:
        return False
    for p in shape.text_frame.paragraphs:
        for r in p.runs:
            if r.text == old_text:
                r.text = new_text
                return True
    # Fallback: buscar sub-string
    for p in shape.text_frame.paragraphs:
        for r in p.runs:
            if old_text in r.text:
                r.text = r.text.replace(old_text, new_text)
                return True
    return False


def replace_shape_full_text(shape, new_text: str) -> bool:
    """Reemplaza todo el texto de un shape que tiene 1 run (simple)."""
    if not shape.has_text_frame:
        return False
    p = shape.text_frame.paragraphs[0]
    if not p.runs:
        return False
    p.runs[0].text = new_text
    # Vaciar runs subsiguientes en el primer parrafo
    for r in p.runs[1:]:
        r.text = ""
    # Vaciar parrafos adicionales
    for extra_p in shape.text_frame.paragraphs[1:]:
        for r in extra_p.runs:
            r.text = ""
    return True


def set_table_cell(table, row, col, new_text: str, style_from=None) -> bool:
    """Escribe texto en cell preservando estilos del primer run.
    Si el cell esta vacio y hay style_from=(row,col), copia el estilo de ese cell."""
    cell = table.cell(row, col)
    tf = cell.text_frame
    if not tf.paragraphs or not tf.paragraphs[0].runs:
        # Cell vacio - crear run, copiando estilo de style_from si se proveyo
        run = tf.paragraphs[0].add_run()
        run.text = new_text
        if style_from:
            src = table.cell(*style_from).text_frame.paragraphs[0].runs[0]
            run.font.size = src.font.size
            run.font.bold = src.font.bold
            run.font.italic = src.font.italic
            if src.font.name:
                run.font.name = src.font.name
            try:
                if src.font.color and src.font.color.rgb:
                    run.font.color.rgb = src.font.color.rgb
            except Exception:
                pass
        return True
    p = tf.paragraphs[0]
    p.runs[0].text = new_text
    for r in p.runs[1:]:
        r.text = ""
    for extra_p in tf.paragraphs[1:]:
        for r in extra_p.runs:
            r.text = ""
    return True


def main():
    prs = Presentation(str(INPUT))
    changes = []
    err = []

    def log(where, ok):
        (changes if ok else err).append(where)

    # ===================== SLIDE 1 (PORTADA) =====================
    s1 = prs.slides[0]

    # Rentabilidades (formato "BIG / ACWI" con salto de línea probable)
    # Shape [5]: "1.0\n3.1" -> "-1.0\n-0.1"
    # Los runs actuales son:
    #  [4] "1 mes"        -> keep
    #  [5] "1.0" | "3.1"   -> "-1.0" | "-0.1"
    # Muchos tienen formato multi-run: primer run BIG, run \n, run ACWI
    RENT_MAP = [
        (5,  ["1.0", "3.1"],  ["-1.0", "-0.1"]),   # 1 mes
        (7,  ["-0.2", "4.2"], ["2.2", "9.0"]),      # 3 meses
        (9,  ["11.9", "15.0"], ["10.4", "13.7"]),    # 3 años
        (11, ["6.8", "7.4"],  ["6.3", "7.1"]),      # 5 años
        (15, ["2.2", "8.1"],  ["1.0", "7.4"]),      # 6 meses
        (17, ["2.0", "7.5"],  ["1.0", "7.4"]),      # YTD
        (19, ["7.5", "19.1"], ["4.8", "15.3"]),     # 1 año
    ]
    for idx, olds, news in RENT_MAP:
        shape = s1.shapes[idx]
        for old, new in zip(olds, news):
            ok = replace_run_text(shape, old, new)
            log(f"S1[{idx}] '{old}'->'{new}'", ok)

    # Inicio [13]: "8, 98" y "9,31" (Mayo) -> "8,73" y "9,19" (Junio)
    # El texto real es fragmentado: "8," + "98" + newline + "9,31"
    shape = s1.shapes[13]
    ok1 = replace_run_text(shape, "98", "73")   # "8, 98" -> "8, 73"
    ok2 = replace_run_text(shape, "9,31", "9,19")
    log(f"S1[13] Inicio", ok1 and ok2)

    # Composición pies (mantener 3 categorias, Cash+Derivados incluido en RF)
    # Junio: 30% RV + 40% RF + 30% Alts = 100%
    # [29] "30% Renta Variable" -> igual (30%)
    # [30] "43 % Renta Fija" -> "40 % Renta Fija"
    ok = replace_run_text(s1.shapes[30], "43", "40")
    log(f"S1[30] RF 43->40", ok)
    # [31] "27 % Alternativos" -> "30 % Alternativos"
    ok = replace_run_text(s1.shapes[31], "27", "30")
    log(f"S1[31] Alts 27->30", ok)

    # Footer titulo Mayo -> Junio
    for idx in [33, 90]:
        ok = replace_run_text(s1.shapes[idx], "Mayo", "Junio")
        log(f"S1[{idx}] Mayo->Junio", ok)

    # Benchmark comparativo - metricas 3 y 5 años (labels)
    BENCH_MAP = [
        (38, "11,9%", "10,4%"),   # Ret 3a BIG
        (39, "6,8%", "6,3%"),     # Ret 5a BIG
        (40, "5,6%", "5,6%"),     # Vol 3a BIG - igual
        (41, "7,2%", "7,2%"),     # Vol 5a BIG - igual
        (44, "15,0%", "13,7%"),   # Ret 3a ACWI
        (45, "7,4%", "7,1%"),     # Ret 5a ACWI
        (46, "8,6%", "8,8%"),     # Vol 3a ACWI
        (47, "9,7%", "9,8%"),     # Vol 5a ACWI
        # MaxDD igual - skip
        (70, "1,35%", "1,12%"),   # Sharpe 3a BIG
        (71, "0,45%", "0,29%"),   # Sharpe 5a BIG
        (72, "1,25%", "1,09%"),   # Sharpe 3a ACWI
        (73, "0,39%", "0,30%"),   # Sharpe 5a ACWI
    ]
    for idx, old, new in BENCH_MAP:
        if old == new:
            continue
        shape = s1.shapes[idx]
        ok = replace_run_text(shape, old, new)
        log(f"S1[{idx}] {old}->{new}", ok)

    # Tabla Rendimiento mensual [22]: agregar Jun 2026 y actualizar YTD Año
    # Row 5 (2026): actualizar cell 6 (Jun) y cell 13 (Año)
    tbl = s1.shapes[22].table
    ok1 = set_table_cell(tbl, 5, 6, "-1.0")   # Jun
    ok2 = set_table_cell(tbl, 5, 13, "1.0")   # Año YTD
    log(f"S1[22] Table 2026 Jun/YTD", ok1 and ok2)

    # ===================== SLIDE 2 (STATS) =====================
    s2 = prs.slides[1]

    # Stats RF BIG [41]: "7.1%" + "4.3" (dur) => "6.8%" + "4.9"
    # [41] text runs: "7.1%", "4.3"
    ok1 = replace_run_text(s2.shapes[41], "7.1%", "6.8%")
    ok2 = replace_run_text(s2.shapes[41], "4.3", "4.9")
    log(f"S2[41] YTW BIG + Dur", ok1 and ok2)
    # [42] Venc BIG: "6.1" -> "7.2"
    ok = replace_run_text(s2.shapes[42], "6.1", "7.2")
    log(f"S2[42] Venc BIG", ok)
    # [60] Rta.Cte BIG: "3.3%" -> "4.1%"
    ok = replace_run_text(s2.shapes[60], "3.3%", "4.1%")
    log(f"S2[60] Rta Cte BIG", ok)

    # Stats RV BIG [43-45]
    ok = replace_run_text(s2.shapes[43], "21.0", "21.0")   # (=)
    ok = replace_run_text(s2.shapes[44], "3.4", "3.3")     # P/B
    log(f"S2[44] P/B BIG 3.4->3.3", ok)
    # P/S [45]: 2.3 (=)

    # Tabla ACWI [22]:
    # R0: "ACWI 60%/40% AGG | 3.9% | 6.2 | 8.4 | 3.0% |  | ACWI 60%/40% AGG | 23.1 | 3.6 | 2.8"
    # Actualizar cells: YTW 3.9->3.8, Venc 8.4 (=), Rta Cte 3.0 (=), P/E 23.1 (=), P/B 3.6->3.7, P/S 2.8->2.9
    tbl2 = s2.shapes[22].table
    ok1 = set_table_cell(tbl2, 0, 1, "3.8%")   # YTW ACWI
    ok2 = set_table_cell(tbl2, 0, 8, "3.7")    # P/B ACWI
    ok3 = set_table_cell(tbl2, 0, 9, "2.9")    # P/S ACWI
    log(f"S2[22] ACWI stats", ok1 and ok2 and ok3)

    # Calidad Crediticia BIG (tabla [22] cells 4)
    # R2 (AAA): 15.8% -> 12.6%
    # R3 (AA): 37.2% -> 37.7%
    # R4 (A): 6.9% -> 6.2%
    # R5 (BBB): 20.7% -> 18.5%
    ok1 = set_table_cell(tbl2, 2, 4, "12.6%")  # AAA
    ok2 = set_table_cell(tbl2, 3, 4, "37.7%")  # AA
    ok3 = set_table_cell(tbl2, 4, 4, "6.2%")   # A
    ok4 = set_table_cell(tbl2, 5, 4, "18.5%")  # BBB
    log(f"S2[22] Calidad AAA/AA/A/BBB", ok1 and ok2 and ok3 and ok4)

    # Calidad BIG tabla [19] - BB/B/CCC
    tbl19 = s2.shapes[19].table
    ok1 = set_table_cell(tbl19, 0, 1, "9.2%")   # BB
    ok2 = set_table_cell(tbl19, 1, 1, "3.3%")   # B
    ok3 = set_table_cell(tbl19, 2, 1, "6.3%")   # CCC
    log(f"S2[19] Calidad BB/B/CCC", ok1 and ok2 and ok3)

    # Footer IG/HY [20]: "Investment Grade: 80.6% | High Yield: 18.9%"
    ok1 = replace_run_text(s2.shapes[20], "80.6%", "74.9%")
    ok2 = replace_run_text(s2.shapes[20], "18.9%", "18.9%")   # (=)
    log(f"S2[20] IG footer", ok1)

    # Sectores RV BIG (tabla [22] cells col 9)
    # R2 Tech: 21.0% -> 20.3%
    # R3 Industrials: 14.0% -> 16.3%
    # R4 Financial: 12.0% -> 13.3%
    # R5 Others (agrupado): 13.0% -> 13.0% (=) - podemos ajustar despues
    ok1 = set_table_cell(tbl2, 2, 9, "20.3%")  # Tech
    ok2 = set_table_cell(tbl2, 3, 9, "16.3%")  # Industrials
    ok3 = set_table_cell(tbl2, 4, 9, "13.3%")  # Financial
    log(f"S2[22] Sectores Tech/Ind/Fin", ok1 and ok2 and ok3)

    # Tabla [23] - Cons Cyc / Health / Comm Serv
    tbl23 = s2.shapes[23].table
    ok1 = set_table_cell(tbl23, 0, 1, "10.9%")  # Cons Cyc 11.0 -> 10.9
    ok2 = set_table_cell(tbl23, 1, 1, "8.9%")   # Health 9.0 -> 8.9
    ok3 = set_table_cell(tbl23, 2, 1, "8.0%")   # Comm Serv 7.0 -> 8.0
    log(f"S2[23] Sectores Cons Cyc/Health/CS", ok1 and ok2 and ok3)

    # Sectores individuales [26-31] - Basic Mat, Cons Def, Real Estate
    ok = replace_run_text(s2.shapes[26], "6.0%", "5.4%")   # Basic Mat 6.0->5.4
    log(f"S2[26] Basic Mat", ok)
    ok = replace_run_text(s2.shapes[28], "5.0%", "5.1%")   # Cons Def 5.0->5.1
    log(f"S2[28] Cons Def", ok)
    ok = replace_run_text(s2.shapes[30], "1.0%", "1.9%")   # Real Estate 1.0->1.9
    log(f"S2[30] Real Estate", ok)

    # Currency exposure [77-88] (números)
    ok = replace_run_text(s2.shapes[77], "74.5%", "72.5%")   # USD
    log(f"S2[77] USD", ok)
    ok = replace_run_text(s2.shapes[79], "6.4%", "7.9%")     # EUR
    log(f"S2[79] EUR", ok)
    ok = replace_run_text(s2.shapes[81], "3.2%", "3.9%")     # GBP
    log(f"S2[81] GBP", ok)
    ok = replace_run_text(s2.shapes[83], "3.2%", "2.7%")     # ORO
    log(f"S2[83] ORO", ok)
    ok = replace_run_text(s2.shapes[85], "3.3%", "2.7%")     # BTC
    log(f"S2[85] BTC", ok)
    ok = replace_run_text(s2.shapes[87], "9.4%", "10.3%")    # Otros
    log(f"S2[87] Otros", ok)

    # Alts [62-64]
    ok = replace_run_text(s2.shapes[62], "13,6%", "16,7%")   # PE
    log(f"S2[62] PE", ok)
    ok = replace_run_text(s2.shapes[63], "6,9%", "7,1%")     # PC
    log(f"S2[63] PC", ok)
    # Alts total 27 -> 29 en shape [64]
    ok1 = replace_run_text(s2.shapes[64], "27,0%", "30,0%")  # Total 27->30
    ok2 = replace_run_text(s2.shapes[64], "3,5%", "2,7%")    # BTC
    ok3 = replace_run_text(s2.shapes[64], "3,2%", "2,7%")    # Oro
    log(f"S2[64] Alts total + BTC + Oro", ok1 and ok2 and ok3)

    # Regional RV [66-68]
    # [66] contiene mucho texto: "Exposición regional renta variable (%) Japan 2.0 % Others 7,0 % UK 6,0% LatAm 14,0%"
    ok = replace_run_text(s2.shapes[66], "2.0", "5.0")   # Japan 2->5
    log(f"S2[66] Japan 2->5", ok)
    ok = replace_run_text(s2.shapes[66], "7,0", "3,0")   # Others 7->3
    log(f"S2[66] Others 7->3", ok)
    ok = replace_run_text(s2.shapes[66], "14,0", "12,0")  # LatAm 14->12
    log(f"S2[66] LatAm 14->12", ok)

    # [67] N America
    ok = replace_run_text(s2.shapes[67], "58,0%", "66,0%")
    log(f"S2[67] N America 58->66", ok)
    # [68] Europe
    ok = replace_run_text(s2.shapes[68], "13,0%", "14,0%")
    log(f"S2[68] Europe 13->14", ok)

    # Footer titulo Mayo -> Junio (S2)
    for idx in [36, 93]:
        ok = replace_run_text(s2.shapes[idx], "Mayo", "Junio")
        log(f"S2[{idx}] Mayo->Junio", ok)

    # ===================== SLIDE 3 (DISCLAIMER) =====================
    s3 = prs.slides[2]
    for idx in [4, 9]:
        ok = replace_run_text(s3.shapes[idx], "Mayo", "Junio")
        log(f"S3[{idx}] Mayo->Junio", ok)

    # ===================== GUARDAR =====================
    prs.save(str(INPUT))

    print(f"\n=== RESUMEN ===")
    print(f"Cambios OK: {len(changes)}")
    print(f"Cambios FALLIDOS: {len(err)}")
    if err:
        print(f"\nFallidos:")
        for e in err:
            print(f"  ! {e}")
    print(f"\nGuardado: {INPUT.name}")


if __name__ == "__main__":
    main()
