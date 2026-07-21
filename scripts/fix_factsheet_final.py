"""
fix_factsheet_final.py
=======================
Fixes finales al Factsheet Junio 2026:
1. Alineacion vertical del cell "-1.0" (Jun 2026) en tabla rendimiento mensual
2. Pie chart Alts: labels con valores 30-Jun + agregar Infraestructura
"""
from pptx import Presentation
from pptx.util import Emu
from pptx.enum.text import PP_ALIGN

FS = 'C:/Users/lmonp/Dropbox/Banca Privada (1)/AMC PAMPA CAPITAL/BIG Factsheets/2026/Pampa Capital AM - Factsheet Junio26.pptx'

def replace_run(shape, old, new):
    if not shape.has_text_frame:
        return False
    for para in shape.text_frame.paragraphs:
        for r in para.runs:
            if old in r.text:
                r.text = r.text.replace(old, new)
                return True
    return False


def main():
    p = Presentation(FS)

    # ============ FIX 1: Alineacion vertical -1.0 en tabla mensual ============
    tbl = p.slides[0].shapes[22].table
    # Copiar props del cell C5 (May) al C6 (Jun)
    src_cell = tbl.cell(5, 5)
    dst_cell = tbl.cell(5, 6)
    src_para = src_cell.text_frame.paragraphs[0]

    dst_cell.margin_top = src_cell.margin_top
    dst_cell.margin_bottom = src_cell.margin_bottom

    dst_para = dst_cell.text_frame.paragraphs[0]
    dst_para.alignment = src_para.alignment
    dst_para.line_spacing = src_para.line_spacing
    dst_para.space_before = src_para.space_before
    dst_para.space_after = src_para.space_after

    print(f"[FIX 1] Cell C6 (Jun 2026) props sincronizados con C5 (May 2026)")

    # ============ FIX 2: Pie Alts labels + agregar Infra ============
    s2 = p.slides[1]

    # Actualizar valores al 30-Jun (matchea Pershing)
    r1 = replace_run(s2.shapes[62], '16,7%', '12,4%')
    print(f"[FIX 2a] PE 16,7% -> 12,4%: {r1}")

    r2 = replace_run(s2.shapes[63], '7,9%', '7,1%')
    print(f"[FIX 2b] PC 7,9% -> 7,1%: {r2}")

    r3 = replace_run(s2.shapes[64], '(30,0%)', '(29,2%)')
    print(f"[FIX 2c] Alts total 30,0% -> 29,2%: {r3}")

    # Agregar Infraestructura al texto [64]
    # Actualmente: 'Activos alternativos (29,2%)\nBitcoin 2,7%\nOro 2,7%'
    # Objetivo: agregar 'Infraestructura 4,3%' como 4ta linea
    tf = s2.shapes[64].text_frame
    src_para = tf.paragraphs[1]  # Bitcoin - copiar formato de aca
    src_run = src_para.runs[0]

    # Agregar nuevo parrafo con Infraestructura
    from pptx.oxml.ns import qn
    from copy import deepcopy

    # Copiar el XML de la ultima linea (Oro) y crear una nueva
    p_oro_xml = tf.paragraphs[2]._p  # ultimo parrafo (Oro)
    new_p_xml = deepcopy(p_oro_xml)

    # Agregar el nuevo parrafo al text frame
    tf._txBody.append(new_p_xml)

    # Refrescar y modificar el texto del nuevo parrafo
    all_paras = list(tf.paragraphs)
    new_para = all_paras[-1]

    # Reemplazar 'Oro' + '2,7%' con 'Infraestructura' + '4,3%'
    for run in new_para.runs:
        if run.text == 'Oro':
            run.text = 'Infraestructura'
        elif '2,7' in run.text:
            run.text = run.text.replace('2,7', '4,3')

    print(f"[FIX 2d] Agregado 'Infraestructura 4,3%' al pie Alts")

    # Guardar
    p.save(FS)
    print(f"\nGuardado: {FS.split(chr(47))[-1]}")

    # Verify
    p2 = Presentation(FS)
    tbl2 = p2.slides[0].shapes[22].table
    c6 = tbl2.cell(5, 6)
    print(f"\nVerify C6: text={c6.text!r} margin_top={c6.margin_top} alignment={c6.text_frame.paragraphs[0].alignment}")

    s2b = p2.slides[1]
    print(f"Verify [62]: {s2b.shapes[62].text_frame.text!r}")
    print(f"Verify [63]: {s2b.shapes[63].text_frame.text!r}")
    print(f"Verify [64]: {s2b.shapes[64].text_frame.text!r}")


if __name__ == "__main__":
    main()
