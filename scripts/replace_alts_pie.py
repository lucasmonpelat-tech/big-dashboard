"""
replace_alts_pie.py
====================
Reemplaza el pie chart de Alts en el Factsheet Junio 2026 (S2 GROUP [61])
con un doughnut chart nativo que incluya:
- Private Equity 12,4%
- Infraestructura 4,3% (NUEVO)
- Private Credit 7,1%
- Bitcoin 2,7%
- Oro 2,7%
Total: 29,2%
"""
from pathlib import Path
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.util import Emu
from pptx.dml.color import RGBColor
from lxml import etree

FS = Path('C:/Users/lmonp/Dropbox/Banca Privada (1)/AMC PAMPA CAPITAL/BIG Factsheets/2026/Pampa Capital AM - Factsheet Junio26.pptx')


def main():
    p = Presentation(str(FS))
    s2 = p.slides[1]

    # 1) Guardar posicion y tamano del GROUP [61] actual
    old_group = s2.shapes[61]
    left = old_group.left
    top = old_group.top
    width = old_group.width
    height = old_group.height
    print(f"GROUP [61] pos=({left},{top}) size=({width},{height})")

    # 2) Eliminar el GROUP viejo (removiendo el XML element del slide)
    sp_element = old_group._element
    sp_element.getparent().remove(sp_element)
    print("GROUP viejo eliminado.")

    # 3) Crear el doughnut chart nativo
    chart_data = CategoryChartData()
    chart_data.categories = ['Private Equity', 'Infraestructura', 'Private Credit', 'Bitcoin', 'Oro']
    chart_data.add_series('Alts', (12.4, 4.3, 7.1, 2.7, 2.7))

    # Insertar chart en la misma pos/size
    graphic_frame = s2.shapes.add_chart(
        XL_CHART_TYPE.DOUGHNUT,
        left, top, width, height,
        chart_data,
    )
    chart = graphic_frame.chart

    # 4) Configurar el look:
    # - Sin titulo, sin leyenda (ya hay text labels alrededor)
    chart.has_title = False
    chart.has_legend = False

    # Colores por slice (paleta Navy + Gold consistente con branding BIG)
    COLORS = [
        RGBColor(0x14, 0x2B, 0x45),  # Private Equity - Navy dark
        RGBColor(0x4A, 0x6C, 0x7C),  # Infraestructura - Blue-gray
        RGBColor(0xA7, 0x9A, 0x8A),  # Private Credit - Warm gray
        RGBColor(0xF7, 0x93, 0x1A),  # Bitcoin - Orange
        RGBColor(0xC9, 0xA2, 0x4A),  # Oro - Gold
    ]

    # Aplicar colores a cada slice
    plot = chart.plots[0]
    for i, point in enumerate(plot.series[0].points):
        fill = point.format.fill
        fill.solid()
        fill.fore_color.rgb = COLORS[i]
        # Sin borde
        line = point.format.line
        line.fill.solid()
        line.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)  # borde blanco
        line.width = Emu(6350)  # 0.5pt

    # Sin data labels (ya estan en text boxes alrededor)
    # Ajustar hole size (grosor del doughnut) - default es ~50%
    try:
        # Set doughnut hole size to 60% for thin ring
        chartSpace = chart._chartSpace
        ns = {'c': 'http://schemas.openxmlformats.org/drawingml/2006/chart'}
        doughnutChart = chartSpace.find('.//c:doughnutChart', ns)
        if doughnutChart is not None:
            firstSliceAng = doughnutChart.find('c:firstSliceAng', ns)
            if firstSliceAng is None:
                firstSliceAng = etree.SubElement(doughnutChart, '{%s}firstSliceAng' % ns['c'])
            firstSliceAng.set('val', '0')

            holeSize = doughnutChart.find('c:holeSize', ns)
            if holeSize is None:
                holeSize = etree.SubElement(doughnutChart, '{%s}holeSize' % ns['c'])
            holeSize.set('val', '40')  # 40% hole - donut delgado
    except Exception as e:
        print(f"WARN hole size: {e}")

    p.save(str(FS))
    print(f"\nGuardado: {FS.name}")
    print(f"5 slices creadas con colores navy/blue-gray/gray/orange/gold")


if __name__ == "__main__":
    main()
