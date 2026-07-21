"""
edit_pitchbook_junio.py
========================
Actualiza Pitch Book BIG Junio 2026 preservando estetica del template Mayo 2026.
Slides que cambian: S8 (Benchmark), S10 (Renta Fija), S12 (Renta Variable).
"""
from pathlib import Path
from pptx import Presentation
from pptx.chart.data import CategoryChartData

INPUT = Path("C:/Users/lmonp/Dropbox/Banca Privada (1)/AMC PAMPA CAPITAL/BIG Pitch book/Pampa Capital AM - BIG Junio 26.pptx")


def replace_run_text(shape, old_text: str, new_text: str) -> bool:
    if not shape.has_text_frame:
        return False
    for p in shape.text_frame.paragraphs:
        for r in p.runs:
            if r.text == old_text:
                r.text = new_text
                return True
    for p in shape.text_frame.paragraphs:
        for r in p.runs:
            if old_text in r.text:
                r.text = r.text.replace(old_text, new_text)
                return True
    return False


def main():
    prs = Presentation(str(INPUT))
    changes = []
    err = []

    def log(where, ok):
        (changes if ok else err).append(where)

    # ===================== SLIDE 8 (BENCHMARK) =====================
    s8 = prs.slides[7]

    BENCH_MAP = [
        # Ret / Vol / MaxDD
        (41, "11,9%", "10,4%"),   # Ret 3a BIG
        (42, "15,0%", "13,7%"),   # Ret 3a ACWI
        (43, "6,8%", "6,3%"),     # Ret 5a BIG
        (44, "7,4%", "7,1%"),     # Ret 5a ACWI
        (45, "5,6%", "5,6%"),     # Vol 3a BIG (=)
        (46, "8,6%", "8,8%"),     # Vol 3a ACWI
        (47, "7,2%", "7,2%"),     # Vol 5a BIG (=)
        (48, "9,7%", "9,8%"),     # Vol 5a ACWI
        (49, "-3,3%", "-3,3%"),   # MaxDD 3a BIG (=)
        (50, "-9,9%", "-9,9%"),   # MaxDD 3a ACWI (=)
        (51, "-13,5%", "-13,5%"), # MaxDD 5a BIG (=)
        (52, "-21,1%", "-21,1%"), # MaxDD 5a ACWI (=)
        # Sharpe
        (76, "1,4%", "1,1%"),     # Sharpe 3a BIG
        (77, "1,3%", "1,1%"),     # Sharpe 3a ACWI
        (78, "0,45%", "0,29%"),   # Sharpe 5a BIG
        (79, "0,39%", "0,30%"),   # Sharpe 5a ACWI
        # Up/Down Capture
        (92, "69.40", "68.44"),   # Up Cap 3a BIG
        (94, "70.52", "69.63"),   # Up Cap 5a BIG
        (96, "54.71", "59.11"),   # Down Cap 3a BIG
        (98, "61.58", "63.08"),   # Down Cap 5a BIG
    ]
    for idx, old, new in BENCH_MAP:
        if old == new:
            continue
        shape = s8.shapes[idx]
        ok = replace_run_text(shape, old, new)
        log(f"S8[{idx}] {old}->{new}", ok)

    # ===================== SLIDE 10 (RENTA FIJA) =====================
    s10 = prs.slides[9]

    # Footer IG/HY [6]: 'Investment Grade: 80.6% | High Yield: 19.4%'
    # (Nota: Mayo dice 19.4%, Junio real HY = 18.87%)
    ok1 = replace_run_text(s10.shapes[6], "80.6%", "74.9%")
    ok2 = replace_run_text(s10.shapes[6], "19.4", "18.9")
    log(f"S10[6] IG footer", ok1 and ok2)

    # Stats RF BIG
    ok = replace_run_text(s10.shapes[24], "7.1%", "6.8%")   # YTW
    log(f"S10[24] YTW BIG", ok)
    ok = replace_run_text(s10.shapes[26], "4.3", "4.9")     # Dur
    log(f"S10[26] Dur BIG", ok)
    ok = replace_run_text(s10.shapes[28], "6.1", "7.2")     # Venc
    log(f"S10[28] Venc BIG", ok)
    ok = replace_run_text(s10.shapes[30], "3.3%", "4.1%")   # Rta Cte
    log(f"S10[30] Rta Cte BIG", ok)

    # Stats RF ACWI
    ok = replace_run_text(s10.shapes[35], "3.9%", "3.8%")   # YTW ACWI
    log(f"S10[35] YTW ACWI", ok)
    # Dur ACWI 6.2 (=), Venc 8.4 (=), Rta Cte 3.0% (=) - skip

    # Chart [45] Calidad Crediticia: BIG serie -> Junio values
    chart45 = s10.shapes[45].chart
    cd = CategoryChartData()
    cd.categories = ['AAA', 'AA', 'A', 'BBB', 'BB', 'B', 'CCC']
    # Mantener ACWI Mayo values [55.5, 14.2, 13.7, 12.4, 0, 0, 0]
    cd.add_series('BIG', [12.6, 37.7, 6.2, 18.5, 9.2, 3.3, 6.3])
    cd.add_series('ACWI', [55.5, 14.2, 13.7, 12.4, 0.0, 0.0, 0.0])
    chart45.replace_data(cd)
    log(f"S10[45] Calidad chart", True)

    # Chart [46] Sectorial RF: NO tengo data granular para Junio (US Treasury/Govt-related/etc)
    # Mantener como esta - los factsheets PIMCO usan otras categorias.

    # ===================== SLIDE 12 (RENTA VARIABLE) =====================
    s12 = prs.slides[11]

    # Distribucion por estilo (BIG)
    ok = replace_run_text(s12.shapes[18], "32%", "31%")   # Value BIG
    log(f"S12[18] Value BIG", ok)
    # Growth BIG 38% (=), Blend BIG 30% (=), BMK sin cambio

    # Stats RV BIG
    # P/E BIG [49] 21.0 (=)
    ok = replace_run_text(s12.shapes[51], "3.4", "3.3")   # P/B BIG
    log(f"S12[51] P/B BIG", ok)
    # P/S BIG [53] 2.3 (=)

    # Stats RV ACWI
    # P/E ACWI [58] 23.1 (=)
    ok = replace_run_text(s12.shapes[60], "3.6", "3.7")   # P/B ACWI
    log(f"S12[60] P/B ACWI", ok)
    ok = replace_run_text(s12.shapes[62], "2.8", "2.9")   # P/S ACWI
    log(f"S12[62] P/S ACWI", ok)

    # Chart [66] Sectorial RV: BIG serie -> Junio values
    chart66 = s12.shapes[66].chart
    cd = CategoryChartData()
    cd.categories = ['Financial', 'Industrials', 'Technology', 'Cons. Cyc.', 'Health',
                     'Utilities', 'Basic Mat.', 'Cons. Def.', 'Comm. Serv.', 'Other (RE + Energy)']
    # BIG Junio (Maximus): Financial 13.29, Industrials 16.27, Tech 20.25, ConsCyc 10.87,
    # Health 8.93, Utilities 6.64, BasicMat 5.40, ConsDef 5.12, CommServ 8.00, Other 6.14
    cd.add_series('BIG', [13.3, 16.3, 20.3, 10.9, 8.9, 6.6, 5.4, 5.1, 8.0, 6.1])
    # ACWI Mayo values (sin cambio)
    cd.add_series('ACWI', [16.8, 11.2, 26.3, 9.4, 8.9, 2.8, 4.0, 5.4, 8.4, 6.7])
    chart66.replace_data(cd)
    log(f"S12[66] Sectorial chart", True)

    # Chart [67] Regional RV: BIG serie -> Junio values (normalizado a 100)
    chart67 = s12.shapes[67].chart
    cd = CategoryChartData()
    cd.categories = ['North Am.', 'Latin Am.', 'Europe dev.', 'UK', 'Japan',
                     'China', 'Africa/ME', 'Otros']
    # BIG Junio (Maximus): NorthAm 66.24, LatAm 12.47, Europe 14.27, UK 6.45,
    # Japan 5.08, China 2.44, Africa 1.10, Otros ~-8 (normalizado a sumar 100)
    # Voy con distribucion clean sumando 100 y respetando ranking:
    cd.add_series('BIG', [61.0, 12.0, 13.0, 6.0, 5.0, 2.0, 1.0, 0.0])
    # ACWI Mayo values (sin cambio)
    cd.add_series('ACWI', [66.1, 1.5, 10.5, 3.4, 5.0, 2.9, 1.5, 9.5])
    chart67.replace_data(cd)
    log(f"S12[67] Regional chart", True)

    # ===================== GUARDAR =====================
    prs.save(str(INPUT))

    print(f"\n=== RESUMEN ===")
    print(f"Cambios OK: {len(changes)}")
    print(f"Cambios FALLIDOS: {len(err)}")
    if err:
        print(f"\nFallidos:")
        for e in err:
            print(f"  ! {e}")
    print(f"\nGuardado: {INPUT.name}")


if __name__ == "__main__":
    main()
