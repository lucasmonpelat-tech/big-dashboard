"""
fix_pitchbook_bars.py
=====================
Redimensiona las barras de charts en Slide 8 del Pitch Book (Ret/Vol/MaxDD +
Sharpe) proporcional a los nuevos valores. Los text labels ya estaban OK del
script edit_pitchbook_junio.py, solo faltaba mover las barras visuales.
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Emu

INPUT = Path("C:/Users/lmonp/Dropbox/Banca Privada (1)/AMC PAMPA CAPITAL/BIG Pitch book/Pampa Capital AM - BIG Junio 26.pptx")


def resize_bar_top_grows(shape, ratio: float, baseline_y: int):
    """Redimensiona barra que crece hacia arriba desde baseline_y (bottom fijo).
    ratio = valor_nuevo / valor_viejo.
    """
    new_height = int(shape.height * ratio)
    new_top = baseline_y - new_height
    shape.top = new_top
    shape.height = new_height


def resize_bar_top_fixed(shape, ratio: float):
    """Redimensiona barra que crece hacia abajo (top fijo)."""
    shape.height = int(shape.height * ratio)


def main():
    prs = Presentation(str(INPUT))
    s8 = prs.slides[7]

    # ============ CHART Ret/Vol/MaxDD (shapes 12-23) ============
    # Baseline (bottom) para Ret y Vol = 1862138 EMU (verificado desde Mayo)
    BASELINE_RETVOL = 1862138

    # Mapa: shape_idx -> ratio (Junio/Mayo)
    RATIOS_CHART1 = {
        12: 10.4 / 11.9,    # Ret 3a BIG
        13: 13.7 / 15.0,    # Ret 3a ACWI
        14: 6.3 / 6.8,      # Ret 5a BIG
        15: 7.1 / 7.4,      # Ret 5a ACWI
        16: 5.6 / 5.6,      # Vol 3a BIG (=)
        17: 8.8 / 8.6,      # Vol 3a ACWI
        18: 7.2 / 7.2,      # Vol 5a BIG (=)
        19: 9.8 / 9.7,      # Vol 5a ACWI
        # MaxDD sin cambios - skip 20-23
    }
    for idx, ratio in RATIOS_CHART1.items():
        if ratio == 1.0:
            continue
        sh = s8.shapes[idx]
        old_h = sh.height
        resize_bar_top_grows(sh, ratio, BASELINE_RETVOL)
        print(f"  S8[{idx}] ratio={ratio:.4f} height {old_h} -> {sh.height}")

    # ============ CHART Sharpe (shapes 60-63) ============
    # Baseline (bottom) = 4553300 EMU
    BASELINE_SHARPE = 4553300

    RATIOS_CHART2 = {
        60: 1.12 / 1.35,   # Sharpe 3a BIG
        61: 1.09 / 1.25,   # Sharpe 3a ACWI
        62: 0.29 / 0.45,   # Sharpe 5a BIG
        63: 0.30 / 0.39,   # Sharpe 5a ACWI
    }
    for idx, ratio in RATIOS_CHART2.items():
        sh = s8.shapes[idx]
        old_h = sh.height
        resize_bar_top_grows(sh, ratio, BASELINE_SHARPE)
        print(f"  S8[{idx}] ratio={ratio:.4f} height {old_h} -> {sh.height}")

    prs.save(str(INPUT))
    print(f"\nGuardado: {INPUT.name}")

    # Tambien tengo que mover los text labels hacia abajo (siguen las barras)
    # Los labels [41-52] estan encima de las barras. Voy a moverlos tambien.
    # En realidad los labels tienen su propia posicion. Verify posiciones.
    prs2 = Presentation(str(INPUT))
    s8_2 = prs2.slides[7]
    for idx in [12, 13, 14, 15, 17, 19, 60, 61, 62, 63]:
        sh = s8_2.shapes[idx]
        print(f"  Verify S8[{idx}] top={sh.top} height={sh.height} bottom={sh.top + sh.height}")


if __name__ == "__main__":
    main()
