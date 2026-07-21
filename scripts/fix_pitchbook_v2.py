"""
fix_pitchbook_v2.py
===================
Fixes al Pitch Book Junio 2026:
1. Redimensionar barras del chart Sharpe (S8) proporcional a nuevos valores
2. Redimensionar barras del chart Ret/Vol/MaxDD (S8) proporcional
3. Actualizar Calidad Crediticia S10:
   - Chart: BBB absorbe SGCB (2.01) + PIMCO NR (1.11), BB absorbe TGF (3.10)
   - Footer: IG 78.0% / HY 22.0% (Sin Calificacion redistribuido a IG)
4. Los mismos cambios al Factsheet Slide 2
"""
from pathlib import Path
from pptx import Presentation
from pptx.chart.data import CategoryChartData

PB = Path("C:/Users/lmonp/Dropbox/Banca Privada (1)/AMC PAMPA CAPITAL/BIG Pitch book/Pampa Capital AM - BIG Junio 26.pptx")
FS = Path("C:/Users/lmonp/Dropbox/Banca Privada (1)/AMC PAMPA CAPITAL/BIG Factsheets/2026/Pampa Capital AM - Factsheet Junio26.pptx")


def replace_run_text(shape, old_text: str, new_text: str) -> bool:
    if not shape.has_text_frame:
        return False
    for p in shape.text_frame.paragraphs:
        for r in p.runs:
            if r.text == old_text:
                r.text = new_text
                return True
    for p in shape.text_frame.paragraphs:
        for r in p.runs:
            if old_text in r.text:
                r.text = r.text.replace(old_text, new_text)
                return True
    return False


def set_table_cell(table, row, col, new_text: str) -> bool:
    cell = table.cell(row, col)
    tf = cell.text_frame
    if not tf.paragraphs or not tf.paragraphs[0].runs:
        run = tf.paragraphs[0].add_run()
        run.text = new_text
        return True
    p = tf.paragraphs[0]
    p.runs[0].text = new_text
    for r in p.runs[1:]:
        r.text = ""
    for extra_p in tf.paragraphs[1:]:
        for r in extra_p.runs:
            r.text = ""
    return True


def resize_bar_bottom_fixed(shape, ratio: float, baseline_y: int):
    """Redimensiona barra positiva (crece hacia arriba desde baseline)."""
    new_height = int(shape.height * ratio)
    shape.top = baseline_y - new_height
    shape.height = new_height


# ============================================================
#                    PITCH BOOK
# ============================================================
def fix_pitchbook():
    prs = Presentation(str(PB))
    s8 = prs.slides[7]
    s10 = prs.slides[9]

    # --- Chart Ret/Vol/MaxDD (S8 shapes 12-19, baseline 1862138) ---
    BASELINE_1 = 1862138
    RATIOS_1 = {
        12: 10.4 / 11.9,   # Ret 3a BIG
        13: 13.7 / 15.0,   # Ret 3a ACWI
        14: 6.3 / 6.8,     # Ret 5a BIG
        15: 7.1 / 7.4,     # Ret 5a ACWI
        17: 8.8 / 8.6,     # Vol 3a ACWI
        19: 9.8 / 9.7,     # Vol 5a ACWI
    }
    for idx, ratio in RATIOS_1.items():
        resize_bar_bottom_fixed(s8.shapes[idx], ratio, BASELINE_1)
        print(f"  S8 chart1 [{idx}] ratio={ratio:.3f}")

    # --- Chart Sharpe (S8 shapes 60-63, baseline 4553300) ---
    BASELINE_2 = 4553300
    RATIOS_2 = {
        60: 1.12 / 1.35,   # Sharpe 3a BIG
        61: 1.09 / 1.25,   # Sharpe 3a ACWI
        62: 0.29 / 0.45,   # Sharpe 5a BIG
        63: 0.30 / 0.39,   # Sharpe 5a ACWI
    }
    for idx, ratio in RATIOS_2.items():
        resize_bar_bottom_fixed(s8.shapes[idx], ratio, BASELINE_2)
        print(f"  S8 chart2 [{idx}] ratio={ratio:.3f}")

    # --- S10 footer IG/HY: 74.9% -> 78.0%, 18.9% -> 22.0% ---
    # Shape [6]: 'Investment Grade: 80.6% | High Yield: 19.4% | Promedio: A-'
    # Previamente cambiamos a 74.9% / 18.9%. Ahora a 78.0% / 22.0%.
    ok1 = replace_run_text(s10.shapes[6], "74.9%", "78.0%")
    ok2 = replace_run_text(s10.shapes[6], "18.9", "22.0")
    print(f"  S10[6] IG/HY footer: IG {ok1}, HY {ok2}")

    # --- S10 chart Calidad Crediticia [45] ---
    # Reclasificacion (SGCB=IG->BBB, TGF=HY->BB, NR restante 1.1% -> BBB tambien)
    # BBB = 18.52 + 2.01 (SGCB) + 1.11 (NR resto) = 21.64
    # BB = 9.22 + 3.10 (TGF) = 12.32
    chart45 = s10.shapes[45].chart
    cd = CategoryChartData()
    cd.categories = ['AAA', 'AA', 'A', 'BBB', 'BB', 'B', 'CCC']
    cd.add_series('BIG', [12.6, 37.7, 6.2, 21.6, 12.3, 3.3, 6.3])
    cd.add_series('ACWI', [55.5, 14.2, 13.7, 12.4, 0.0, 0.0, 0.0])
    chart45.replace_data(cd)
    print(f"  S10[45] Calidad chart actualizado con reclasificacion")

    prs.save(str(PB))
    print(f"\nGuardado: {PB.name}")


# ============================================================
#                    FACTSHEET
# ============================================================
def fix_factsheet():
    prs = Presentation(str(FS))
    s2 = prs.slides[1]

    # Footer IG/HY [20] Factsheet: mismo cambio
    ok1 = replace_run_text(s2.shapes[20], "74.9%", "78.0%")
    # HY 18.9% ya estaba correcto en el Factsheet (18.87 redondeado), lo dejamos
    # Pero si Lucas quiere HY 22.0%, cambiamos:
    ok2 = replace_run_text(s2.shapes[20], "18.9%", "22.0%")
    print(f"  S2[20] Footer IG/HY: IG {ok1}, HY {ok2}")

    # Tabla Calidad Crediticia [22] cells 4 (AAA-BBB) y tabla [19] (BB-CCC)
    tbl2 = s2.shapes[22].table
    tbl19 = s2.shapes[19].table

    # BBB: 18.5 -> 21.6
    set_table_cell(tbl2, 5, 4, "21.6%")
    # BB: 9.2 -> 12.3
    set_table_cell(tbl19, 0, 1, "12.3%")
    print(f"  S2 Calidad crediticia BBB/BB actualizados")

    prs.save(str(FS))
    print(f"Guardado: {FS.name}")


def main():
    print("=== FIX PITCH BOOK ===")
    fix_pitchbook()
    print("\n=== FIX FACTSHEET ===")
    fix_factsheet()


if __name__ == "__main__":
    main()
